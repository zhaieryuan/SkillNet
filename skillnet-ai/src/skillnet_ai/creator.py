import os
import json
import re
import ast
import logging
from typing import List, Optional, Dict, Any

import requests
from openai import OpenAI
from skillnet_ai.prompts import (
    CANDIDATE_METADATA_SYSTEM_PROMPT,
    CANDIDATE_METADATA_USER_PROMPT_TEMPLATE,
    SKILL_CONTENT_SYSTEM_PROMPT,
    SKILL_CONTENT_USER_PROMPT_TEMPLATE,
    GITHUB_SKILL_SYSTEM_PROMPT,
    GITHUB_SKILL_USER_PROMPT_TEMPLATE,
    OFFICE_SKILL_SYSTEM_PROMPT,
    OFFICE_SKILL_USER_PROMPT_TEMPLATE,
    PROMPT_SKILL_SYSTEM_PROMPT,
    PROMPT_SKILL_USER_PROMPT_TEMPLATE
)

logger = logging.getLogger(__name__)

class SkillCreator:
    """
    Creates Skill packages from execution trajectories using OpenAI-compatible LLMs.
    """
    
    def __init__(self, api_key: Optional[str] = None, base_url: Optional[str] = None, model: str = "gpt-4o"):
        self.api_key = api_key or os.getenv("API_KEY")
        self.base_url = base_url or os.getenv("BASE_URL") or "https://api.openai.com/v1"
        self.model = model
        
        if not self.api_key:
            raise ValueError("API Key is missing. Please provide it in init or set API_KEY environment variable.")
            
        self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)

    def _get_llm_response(self, messages: List[dict]) -> str:
        """Helper to call LLM and get string content."""
        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=messages
            )
            return response.choices[0].message.content
        except Exception as e:
            logger.error(f"LLM Call Failed: {e}")
            raise

    def create_from_trajectory(self, trajectory: str, output_dir: str = ".") -> List[str]:
        """
        Main entry point: Analyze trajectory and create skill files.
        
        Args:
            trajectory: The string content of the user's action log/trajectory.
            output_dir: The directory where skills should be saved.
            
        Returns:
            List of paths to the created skill directories.
        """
        logger.info("Step 1: Analyzing trajectory to identify skills...")
        
        # 1. Create Metadata
        meta_messages = [
            {"role": "system", "content": CANDIDATE_METADATA_SYSTEM_PROMPT},
            {"role": "user", "content": CANDIDATE_METADATA_USER_PROMPT_TEMPLATE.format(trajectory=trajectory)}
        ]
        
        raw_meta_response = self._get_llm_response(meta_messages)
        candidates = self._parse_candidate_metadata(raw_meta_response)
        
        if not candidates:
            logger.warning("No skills identified in the trajectory.")
            return []

        created_paths = []
        
        # 2. Create Content for each candidate
        for cand in candidates:
            name = cand.get("name")
            description = cand.get("description")
            logger.info(f"Creating content for skill: {name}...")
            
            content_messages = [
                {"role": "system", "content": SKILL_CONTENT_SYSTEM_PROMPT},
                {"role": "user", "content": SKILL_CONTENT_USER_PROMPT_TEMPLATE.format(
                    trajectory=trajectory, name=name, description=description
                )}
            ]
            
            raw_content_response = self._get_llm_response(content_messages)
            
            # 3. Parse and Save Files
            self._save_skill_files(raw_content_response, output_dir)
            created_paths.append(os.path.join(output_dir, name))
            
        return created_paths

    def _parse_candidate_metadata(self, llm_output: str) -> List[dict]:
        """Extract JSON from the LLM output tags."""
        try:
            # Look for content between <Skill_Candidate_Metadata> tags
            if "<Skill_Candidate_Metadata>" in llm_output:
                json_str = llm_output.split("<Skill_Candidate_Metadata>")[1].split("</Skill_Candidate_Metadata>")[0]
            else:
                # Fallback: try to find the first JSON list block
                json_str = llm_output
            
            # clean markdown code blocks if present
            json_str = json_str.replace("```json", "").replace("```", "").strip()
            return json.loads(json_str)
        except Exception as e:
            logger.error(f"Failed to parse metadata JSON: {e}")
            return []

    def _save_skill_files(self, llm_output: str, output_base_dir: str) -> List[str]:
        """Parse the FILE blocks and write them to disk."""
        # Regex to find: ## FILE: path \n ```lang \n content \n ```
        pattern = re.compile(r'##\s*FILE:\s*(.+?)\s*\n```(?:\w*)\n(.*?)```', re.DOTALL)
        matches = pattern.findall(llm_output)
        
        created_files = []
        
        if not matches:
            logger.warning("No file blocks found in LLM output.")
            return created_files

        for file_path, content in matches:
            file_path = file_path.strip()
            full_path = os.path.join(output_base_dir, file_path)
            
            # Create directory if missing
            os.makedirs(os.path.dirname(full_path), exist_ok=True)
            
            try:
                with open(full_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                logger.info(f"Saved: {full_path}")
                created_files.append(full_path)
            except IOError as e:
                logger.error(f"Failed to write {full_path}: {e}")
        
        return created_files

    def create_from_office(
        self,
        file_path: str,
        output_dir: str = "./generated_skills"
    ) -> List[str]:
        """
        Create a skill package from an Office document (PDF, PPT, Word).
        
        Args:
            file_path: Path to the office document
            output_dir: Directory where new skills will be saved
            
        Returns:
            List of paths to created skill directories
        """
        logger.info(f"Creating skill from office document: {file_path}")
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if not _OfficeReader.is_supported(file_path):
            raise ValueError(
                f"Unsupported file type. Supported: {_OfficeReader.SUPPORTED_EXTENSIONS}"
            )
        
        # Extract text content
        document_content = _OfficeReader.extract_text(file_path)
        if not document_content.strip():
            logger.warning("No text content extracted from document")
            return []
        
        filename = os.path.basename(file_path)
        file_type = _OfficeReader.get_file_type(file_path)
        
        # Generate skill using LLM
        user_prompt = OFFICE_SKILL_USER_PROMPT_TEMPLATE.format(
            filename=filename,
            file_type=file_type,
            document_content=document_content
        )
        
        messages = [
            {"role": "system", "content": OFFICE_SKILL_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt}
        ]
        
        try:
            response = self._get_llm_response(messages)
            created_files = self._save_github_skill_files(response, output_dir)
            
            # Extract unique skill directories
            skill_dirs = set()
            for created_file in created_files:
                rel_path = os.path.relpath(created_file, output_dir)
                skill_dir = rel_path.split(os.sep)[0]
                skill_dirs.add(os.path.join(output_dir, skill_dir))
            
            logger.info(f"Skill created from office document: {file_path}")
            return list(skill_dirs)
            
        except Exception as e:
            logger.error(f"Failed to create skill from office document: {e}")
            return []

    def create_from_prompt(
        self,
        user_input: str,
        output_dir: str = "./generated_skills"
    ) -> List[str]:
        """
        Create a skill package from user's direct description.
        
        Args:
            user_input: User's description of the skill to create
            output_dir: Directory where new skills will be saved
            
        Returns:
            List of paths to created skill directories
        """
        logger.info("Creating skill from user prompt")
        
        if not user_input or not user_input.strip():
            raise ValueError("User input cannot be empty")
        
        # Generate skill using LLM
        user_prompt = PROMPT_SKILL_USER_PROMPT_TEMPLATE.format(
            user_input=user_input
        )
        
        messages = [
            {"role": "system", "content": PROMPT_SKILL_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt}
        ]
        
        try:
            response = self._get_llm_response(messages)
            created_files = self._save_github_skill_files(response, output_dir)
            
            # Extract unique skill directories
            skill_dirs = set()
            for created_file in created_files:
                rel_path = os.path.relpath(created_file, output_dir)
                skill_dir = rel_path.split(os.sep)[0]
                skill_dirs.add(os.path.join(output_dir, skill_dir))
            
            logger.info("Skill created from user input")
            return list(skill_dirs)
            
        except Exception as e:
            logger.error(f"Failed to create skill from user input: {e}")
            return []

    def create_from_github(
        self,
        github_url: str,
        output_dir: str = "./generated_skills",
        api_token: Optional[str] = None,
        max_files: int = 50
    ) -> List[str]:
        """Create a skill package from a GitHub repository."""
        logger.info(f"Creating skill from GitHub: {github_url}")

        try:

            fetcher = _GitHubFetcher(api_token=api_token)
            owner, repo, branch, _ = fetcher.parse_github_url(github_url)
            logger.info(f"Parsed: {owner}/{repo} @ {branch}")

            # Fetch repository data
            repo_data = self._fetch_github_repo_data(fetcher, owner, repo, branch, max_files)
            
            if not repo_data:
                logger.error("Failed to fetch repository data")
                return []

            # Generate skill content
            skill_content = self._generate_github_skill_content(repo_data)
            if not skill_content:
                logger.error("Failed to generate skill content")
                return []

            # Save skill package
            skill_name = repo_data["metadata"]["name"].lower().replace(" ", "-").replace("_", "-")
            created_files = self._save_github_skill_files(skill_content, output_dir)
            
            # Extract unique skill directories from created files
            skill_dirs = set()
            for file_path in created_files:
                rel_path = os.path.relpath(file_path, output_dir)
                skill_dir = rel_path.split(os.sep)[0]
                skill_dirs.add(os.path.join(output_dir, skill_dir))
            
            logger.info(f"Skill created successfully from GitHub: {github_url}")
            return list(skill_dirs) if skill_dirs else [os.path.join(output_dir, skill_name)]

        except Exception as e:
            logger.error(f"Failed to create skill from GitHub: {e}")
            return []

    def _fetch_github_repo_data(
        self,
        fetcher: "_GitHubFetcher",
        owner: str,
        repo: str,
        branch: str,
        max_files: int
    ) -> Optional[Dict[str, Any]]:
        """Fetch repository metadata, README, file tree, and code analysis."""
        logger.info("Fetching repository data...")

        metadata = fetcher.fetch_repo_metadata(owner, repo)
        branch = metadata.get("default_branch", branch)
        readme = fetcher.fetch_readme(owner, repo, branch)
        file_tree = fetcher.fetch_file_tree(owner, repo, branch)
        languages = fetcher.fetch_languages(owner, repo)
        code_analysis = self._analyze_github_code_files(
            fetcher, owner, repo, branch, file_tree, max_files
        )

        return {
            "metadata": metadata,
            "readme": readme,
            "file_tree": file_tree,
            "languages": languages,
            "code_analysis": code_analysis,
            "github_url": f"https://github.com/{owner}/{repo}"
        }

    def _analyze_github_code_files(
        self,
        fetcher: "_GitHubFetcher",
        owner: str,
        repo: str,
        branch: str,
        file_tree: List[Dict],
        max_files: int
    ) -> Dict[str, Any]:
        """Analyze code files across supported languages to extract class and function signatures."""
        logger.info("Analyzing code files...")

        # Filter supported code files (Python, JS/TS, Java, Go, C/C++, Rust)
        code_files = [
            f for f in file_tree
            if f.get("type") == "file" and _CodeAnalyzer.is_supported(f.get("path", ""))
        ][:max_files]

        analyzed = []
        for file_info in code_files:
            file_path = file_info.get("path", "")
            content = fetcher.fetch_file_content(owner, repo, file_path, branch)
            
            if content:
                analysis = _CodeAnalyzer.analyze(content, file_path)
                if analysis.get("classes") or analysis.get("functions"):
                    analyzed.append({"file": file_path, **analysis})
                    logger.debug(
                        f"Analyzed {file_path}: {len(analysis.get('classes', []))} classes, "
                        f"{len(analysis.get('functions', []))} functions"
                    )

        total_classes = sum(len(f.get("classes", [])) for f in analyzed)
        total_functions = sum(len(f.get("functions", [])) for f in analyzed)

        logger.info(
            f"Code analysis complete: {len(analyzed)} files, "
            f"{total_classes} classes, {total_functions} functions"
        )

        return {
            "files_analyzed": len(analyzed),
            "total_classes": total_classes,
            "total_functions": total_functions,
            "files": analyzed
        }

    def _generate_github_skill_content(
        self, 
        repo_data: Dict[str, Any],
        max_retries: int = 2
    ) -> Optional[str]:
        """Generate skill content from repository data using LLM."""
        logger.info("Generating skill content with LLM...")

        metadata = repo_data["metadata"]
        code_summary = self._build_code_summary(repo_data.get("code_analysis", {}))
        file_tree_str = self._format_file_tree(repo_data.get("file_tree", [])[:100])

        languages = repo_data.get("languages", {})
        lang_str = ", ".join([f"{lang}: {pct}%" for lang, pct in languages.items()][:5])

        readme_content = repo_data.get("readme") or "No README available"
        readme_truncated = readme_content[:15000]

        user_prompt = GITHUB_SKILL_USER_PROMPT_TEMPLATE.format(
            repo_name=metadata.get("full_name", metadata.get("name", "unknown")),
            repo_url=repo_data.get("github_url", ""),
            repo_description=metadata.get("description") or "No description available",
            language=metadata.get("language") or "Unknown",
            languages_breakdown=lang_str or "N/A",
            stars=metadata.get("stars", 0),
            topics=", ".join(metadata.get("topics", [])) if metadata.get("topics") else "None",
            readme_content=readme_truncated,
            file_tree=file_tree_str,
            code_summary=code_summary
        )

        messages = [
            {"role": "system", "content": GITHUB_SKILL_SYSTEM_PROMPT},
            {"role": "user", "content": user_prompt}
        ]

        # Retry mechanism with content validation
        for attempt in range(max_retries + 1):
            try:
                response = self._get_llm_response(messages)
                
                # Validate response has required content
                if self._validate_skill_content(response):
                    return response
                
                if attempt < max_retries:
                    logger.warning(f"Generated content incomplete, retrying ({attempt + 1}/{max_retries})...")
                else:
                    logger.warning("Generated content may be incomplete, using best result.")
                    return response
                    
            except Exception as e:
                logger.error(f"LLM call failed (attempt {attempt + 1}): {e}")
                if attempt == max_retries:
                    return None
        
        return None

    def _validate_skill_content(self, content: str) -> bool:
        """Check if generated content has required SKILL.md structure."""
        if not content:
            return False
        
        has_skill_md = "SKILL.md" in content
        has_file_block = "## FILE:" in content
        has_frontmatter = "---" in content and "name:" in content
        min_length = len(content) >= 1000
        
        return has_skill_md and has_file_block and has_frontmatter and min_length

    def _build_code_summary(self, code_analysis: Dict[str, Any]) -> str:
        """Build code analysis summary for LLM prompt."""
        if not code_analysis or not code_analysis.get("files"):
            return "No code analysis available."

        # Collect language breakdown from analyzed files
        lang_counts = {}
        for file_data in code_analysis.get("files", []):
            lang = file_data.get("language", "unknown")
            lang_counts[lang] = lang_counts.get(lang, 0) + 1
        lang_info = ", ".join(
            f"{lang}: {count}" for lang, count in
            sorted(lang_counts.items(), key=lambda x: -x[1])
        )

        lines = [
            f"Analyzed {code_analysis.get('files_analyzed', 0)} code files ({lang_info}):",
            f"- Total Classes: {code_analysis.get('total_classes', 0)}",
            f"- Total Functions: {code_analysis.get('total_functions', 0)}",
            "",
            "Key components found:"
        ]

        # Add classes and functions from all analyzed files
        max_lines = 5 + len(code_analysis.get("files", [])) * 4
        for file_data in code_analysis.get("files", []):
            file_path = file_data.get("file", "")
            classes = file_data.get("classes", [])
            functions = file_data.get("functions", [])

            for cls in classes[:3]:
                docstring = (cls.get("docstring") or "")[:100]
                lines.append(f"- Class `{cls['name']}` in {file_path}")
                if docstring:
                    lines.append(f"  {docstring}...")
            
            for func in functions[:3]:
                params = func.get("parameters", [])[:3]
                sig = f"{func['name']}({', '.join(params)})"
                lines.append(f"- Function `{sig}` in {file_path}")

        return "\n".join(lines[:max_lines])

    def _format_file_tree(self, file_tree: List[Dict]) -> str:
        """Format file tree for LLM prompt."""
        if not file_tree:
            return "No file tree available."

        lines = []
        for item in file_tree[:50]:
            path = item.get("path", "")
            item_type = item.get("type", "file")
            icon = "📁" if item_type == "dir" else "📄"
            lines.append(f"{icon} {path}")

        if len(file_tree) > 50:
            lines.append(f"... and {len(file_tree) - 50} more files")

        return "\n".join(lines)

    def _save_github_skill_files(self, llm_output: str, output_base_dir: str) -> List[str]:
        """Parse FILE blocks and write to disk, handling nested code blocks."""
        created_files = []
        parts = re.split(r'##\s*FILE:\s*', llm_output)
        
        if len(parts) < 2:
            logger.warning("No file blocks found in LLM output.")
            return created_files
        
        for part in parts[1:]:
            lines = part.split('\n', 1)
            if len(lines) < 2:
                continue
            
            file_path = lines[0].strip()
            rest = lines[1]
            
            match = re.match(r'```(?:\w*)\n', rest)
            if not match:
                continue
            
            content_start = match.end()
            content = rest[content_start:]
            
            # Find closing ``` by tracking nested code blocks
            in_nested_block = False
            end_pos = -1
            i = 0
            
            while i < len(content):
                # Check for ``` at start of line
                if content[i:i+3] == '```' and (i == 0 or content[i-1] == '\n'):
                    if not in_nested_block:
                        after = content[i+3:i+50].split('\n')[0].strip()
                        if after == '':
                            end_pos = i
                            break
                        else:
                            in_nested_block = True
                    else:
                        in_nested_block = False
                i += 1
            
            if end_pos == -1:
                end_pos = content.rfind('\n```')
                if end_pos == -1:
                    end_pos = len(content)
            
            file_content = content[:end_pos]
            full_path = os.path.join(output_base_dir, file_path.strip())
            
            os.makedirs(os.path.dirname(full_path), exist_ok=True)
            
            try:
                with open(full_path, 'w', encoding='utf-8') as f:
                    f.write(file_content)
                logger.info(f"Saved: {full_path}")
                created_files.append(full_path)
            except IOError as e:
                logger.error(f"Failed to write {full_path}: {e}")
        
        return created_files


class _GitHubFetcher:
    """Fetches content from GitHub repositories via API."""

    EXCLUDED_DIRS = {
        "node_modules", "__pycache__", ".git", ".venv", "venv",
        "env", ".env", "build", "dist", ".pytest_cache",
        ".mypy_cache", "htmlcov", ".tox", ".eggs"
    }

    def __init__(self, api_token: Optional[str] = None):
        self.api_token = api_token or os.getenv("GITHUB_TOKEN")
        self.session = requests.Session()
        self.session.headers.update({
            "Accept": "application/vnd.github.v3+json",
            "User-Agent": "SkillNet-AI/1.0"
        })
        if self.api_token:
            self.session.headers.update({"Authorization": f"token {self.api_token}"})

    def _request_with_retry(
        self, 
        url: str, 
        timeout: int = 10, 
        max_retries: int = 3,
        base_delay: float = 1.0
    ) -> Optional[requests.Response]:
        """HTTP GET with exponential backoff and rate limit handling."""
        import time
        
        for attempt in range(1, max_retries + 1):
            try:
                response = self.session.get(url, timeout=timeout)
                
                if response.status_code == 403:
                    remaining = response.headers.get("X-RateLimit-Remaining", "?")
                    if remaining == "0":
                        reset_time = int(response.headers.get("X-RateLimit-Reset", 0))
                        wait_seconds = max(0, reset_time - int(time.time()))
                        logger.warning(f"GitHub rate limit exceeded. Resets in {wait_seconds}s")
                        if wait_seconds < 60:
                            time.sleep(wait_seconds + 1)
                            continue
                
                return response
                
            except requests.exceptions.Timeout:
                if attempt < max_retries:
                    delay = base_delay * (2 ** (attempt - 1))
                    logger.warning(f"Timeout (attempt {attempt}/{max_retries}), retry in {delay:.1f}s")
                    time.sleep(delay)
                else:
                    logger.error(f"Request failed after {max_retries} attempts: {url}")
                    return None
                    
            except requests.exceptions.ConnectionError:
                if attempt < max_retries:
                    delay = base_delay * (2 ** (attempt - 1))
                    logger.warning(f"Connection error (attempt {attempt}/{max_retries}), retry in {delay:.1f}s")
                    time.sleep(delay)
                else:
                    logger.error(f"Connection failed after {max_retries} attempts: {url}")
                    return None
                    
            except requests.exceptions.RequestException as e:
                logger.error(f"Request failed: {e}")
                return None
        
        return None

    def parse_github_url(self, url: str) -> tuple:
        """Parse GitHub URL to extract owner, repo, branch, and optional path."""
        url = url.rstrip("/")
        if url.endswith(".git"):
            url = url[:-4]

        if "github.com/" in url:
            parts = url.split("github.com/")[-1].split("/")
            if len(parts) < 2:
                raise ValueError(f"Invalid GitHub URL format: {url}")
            
            owner, repo = parts[0], parts[1]
            branch = "main"
            path = ""

            if len(parts) > 3 and parts[2] in ("tree", "blob"):
                branch = parts[3]
                path = "/".join(parts[4:]) if len(parts) > 4 else ""
            
            return owner, repo, branch, path
        
        raise ValueError(f"Invalid GitHub URL: {url}")

    def fetch_repo_metadata(self, owner: str, repo: str) -> Dict[str, Any]:
        """Fetch repository metadata from GitHub API."""
        url = f"https://api.github.com/repos/{owner}/{repo}"
        
        response = self._request_with_retry(url, timeout=10)
        if response is None:
            logger.warning("Failed to fetch repo metadata: request failed")
            return {"name": repo, "full_name": f"{owner}/{repo}"}
        
        try:
            response.raise_for_status()
            data = response.json()

            return {
                "name": data.get("name", repo),
                "full_name": data.get("full_name", f"{owner}/{repo}"),
                "description": data.get("description"),
                "url": data.get("html_url"),
                "homepage": data.get("homepage"),
                "stars": data.get("stargazers_count", 0),
                "forks": data.get("forks_count", 0),
                "language": data.get("language"),
                "topics": data.get("topics", []),
                "license_name": data.get("license", {}).get("name") if data.get("license") else None,
                "default_branch": data.get("default_branch", "main")
            }
        except requests.RequestException as e:
            logger.warning(f"Failed to fetch repo metadata: {e}")
            return {"name": repo, "full_name": f"{owner}/{repo}"}

    def fetch_readme(self, owner: str, repo: str, branch: str = "main") -> Optional[str]:
        """Fetch README content from repository."""
        readme_names = ["README.md", "README.rst", "README.txt", "README"]
        
        for readme_name in readme_names:
            url = f"https://raw.githubusercontent.com/{owner}/{repo}/{branch}/{readme_name}"
            response = self._request_with_retry(url, timeout=10)
            if response and response.status_code == 200:
                logger.info(f"Found README: {readme_name}")
                return response.text
        
        logger.warning("No README found in repository")
        return None

    def fetch_file_tree(self, owner: str, repo: str, branch: str = "main") -> List[Dict]:
        """Fetch repository file tree structure."""
        url = f"https://api.github.com/repos/{owner}/{repo}/git/trees/{branch}?recursive=1"
        
        response = self._request_with_retry(url, timeout=15)
        if response is None:
            logger.warning("Failed to fetch file tree: request failed")
            return []
        
        try:
            response.raise_for_status()
            data = response.json()

            file_tree = []
            for item in data.get("tree", []):
                path = item.get("path", "")
                
                # Skip excluded directories
                if any(excluded in path for excluded in self.EXCLUDED_DIRS):
                    continue
                
                file_tree.append({
                    "path": path,
                    "type": "dir" if item.get("type") == "tree" else "file",
                    "size": item.get("size")
                })
            
            logger.info(f"Fetched file tree: {len(file_tree)} items")
            return file_tree

        except requests.RequestException as e:
            logger.warning(f"Failed to fetch file tree: {e}")
            return []

    def fetch_languages(self, owner: str, repo: str) -> Dict[str, float]:
        """Fetch language breakdown from GitHub API."""
        url = f"https://api.github.com/repos/{owner}/{repo}/languages"
        
        response = self._request_with_retry(url, timeout=10)
        if response is None:
            logger.warning("Failed to fetch languages: request failed")
            return {}
        
        try:
            response.raise_for_status()
            data = response.json()

            if not data:
                return {}

            total_bytes = sum(data.values())
            return {
                lang: round((bytes_count / total_bytes) * 100, 2)
                for lang, bytes_count in data.items()
            }

        except requests.RequestException as e:
            logger.warning(f"Failed to fetch languages: {e}")
            return {}

    def fetch_file_content(
        self, owner: str, repo: str, file_path: str, branch: str = "main"
    ) -> Optional[str]:
        """Fetch content of a specific file."""
        url = f"https://raw.githubusercontent.com/{owner}/{repo}/{branch}/{file_path}"
        
        response = self._request_with_retry(url, timeout=10, max_retries=2)
        if response and response.status_code == 200:
            return response.text
        
        return None


class _CodeAnalyzer:
    """Dispatcher that routes code analysis to language-specific analyzers."""

    EXTENSION_MAP = {
        '.py': 'python',
        '.js': 'javascript', '.jsx': 'javascript', '.mjs': 'javascript',
        '.ts': 'typescript', '.tsx': 'typescript',
        '.java': 'java',
        '.go': 'go',
        '.c': 'c', '.h': 'c',
        '.cpp': 'cpp', '.hpp': 'cpp', '.cc': 'cpp', '.cxx': 'cpp', '.hh': 'cpp',
        '.rs': 'rust',
    }

    SUPPORTED_EXTENSIONS = set(EXTENSION_MAP.keys())

    @staticmethod
    def is_supported(file_path: str) -> bool:
        """Check if file type is supported for code analysis."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in _CodeAnalyzer.SUPPORTED_EXTENSIONS

    @staticmethod
    def get_language(file_path: str) -> Optional[str]:
        """Get language identifier from file path."""
        ext = os.path.splitext(file_path)[1].lower()
        return _CodeAnalyzer.EXTENSION_MAP.get(ext)

    @staticmethod
    def analyze(content: str, file_path: str) -> Dict[str, Any]:
        """Analyze code file and extract structural information."""
        language = _CodeAnalyzer.get_language(file_path)

        if language == 'python':
            result = _PythonCodeAnalyzer.analyze(content, file_path)
        elif language:
            result = _RegexCodeAnalyzer.analyze(content, file_path, language)
        else:
            return {"classes": [], "functions": [], "language": "unknown"}

        result["language"] = language
        return result


class _PythonCodeAnalyzer:
    """Internal class for analyzing Python code using AST."""

    @staticmethod
    def analyze(content: str, file_path: str) -> Dict[str, Any]:
        """Analyze Python file to extract classes and functions."""
        try:
            tree = ast.parse(content)
        except SyntaxError as e:
            logger.debug(f"Syntax error in {file_path}: {e}")
            return {"classes": [], "functions": []}

        classes = []
        functions = []

        for node in ast.iter_child_nodes(tree):
            if isinstance(node, ast.ClassDef):
                classes.append(_PythonCodeAnalyzer._extract_class(node))
            elif isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
                functions.append(_PythonCodeAnalyzer._extract_function(node))

        return {"classes": classes, "functions": functions}

    @staticmethod
    def _extract_class(node: ast.ClassDef) -> Dict[str, Any]:
        """Extract class signature from AST node."""
        base_classes = []
        for base in node.bases:
            if isinstance(base, ast.Name):
                base_classes.append(base.id)
            elif isinstance(base, ast.Attribute):
                if isinstance(base.value, ast.Name):
                    base_classes.append(f"{base.value.id}.{base.attr}")
                else:
                    base_classes.append(base.attr)

        methods = []
        for item in node.body:
            if isinstance(item, (ast.FunctionDef, ast.AsyncFunctionDef)):
                methods.append(_PythonCodeAnalyzer._extract_function(item))

        docstring = ast.get_docstring(node)
        return {
            "name": node.name,
            "base_classes": base_classes,
            "docstring": docstring[:200] if docstring else None,
            "methods": methods
        }

    @staticmethod
    def _extract_function(node) -> Dict[str, Any]:
        """Extract function signature from AST node."""
        params = []
        for arg in node.args.args:
            param_name = arg.arg
            if arg.annotation:
                try:
                    param_name += f": {ast.unparse(arg.annotation)}"
                except Exception:
                    pass
            params.append(param_name)

        return_type = None
        if node.returns:
            try:
                return_type = ast.unparse(node.returns)
            except Exception:
                pass

        decorators = []
        for decorator in node.decorator_list:
            try:
                decorators.append(ast.unparse(decorator))
            except Exception:
                if isinstance(decorator, ast.Name):
                    decorators.append(decorator.id)

        docstring = ast.get_docstring(node)
        return {
            "name": node.name,
            "parameters": params,
            "return_type": return_type,
            "docstring": docstring[:200] if docstring else None,
            "is_async": isinstance(node, ast.AsyncFunctionDef),
            "decorators": decorators
        }


class _RegexCodeAnalyzer:
    """Regex-based code analyzer for extracting structural info from non-Python languages."""

    _KEYWORDS = frozenset({
        'if', 'else', 'for', 'while', 'do', 'switch', 'case', 'default',
        'catch', 'try', 'finally', 'throw', 'return', 'break', 'continue',
        'sizeof', 'typeof', 'instanceof', 'new', 'delete', 'import', 'export',
        'package', 'include', 'define', 'typedef', 'using', 'namespace',
    })

    _CLASS_PATTERNS = {
        'javascript': [
            re.compile(r'(?:export\s+)?class\s+(\w+)(?:\s+extends\s+([\w.]+))?'),
        ],
        'typescript': [
            re.compile(
                r'(?:export\s+)?(?:abstract\s+)?class\s+(\w+)'
                r'(?:\s+extends\s+([\w.]+))?(?:\s+implements\s+([\w.,\s]+))?'
            ),
            re.compile(r'(?:export\s+)?interface\s+(\w+)(?:\s+extends\s+([\w.,\s]+))?'),
        ],
        'java': [
            re.compile(
                r'(?:public|protected|private)?\s*(?:abstract\s+|final\s+|static\s+)*'
                r'class\s+(\w+)(?:\s+extends\s+(\w+))?(?:\s+implements\s+([\w,\s]+))?'
            ),
            re.compile(r'(?:public|protected|private)?\s*interface\s+(\w+)(?:\s+extends\s+([\w,\s]+))?'),
            re.compile(r'(?:public|protected|private)?\s*enum\s+(\w+)'),
        ],
        'go': [
            re.compile(r'type\s+(\w+)\s+struct\b'),
            re.compile(r'type\s+(\w+)\s+interface\b'),
        ],
        'c': [
            re.compile(r'(?:typedef\s+)?struct\s+(\w+)'),
        ],
        'cpp': [
            re.compile(
                r'(?:template\s*<[^>]*>\s*)?class\s+(\w+)'
                r'(?:\s*:\s*(?:public|protected|private)\s+([\w:]+))?'
            ),
            re.compile(r'(?:typedef\s+)?struct\s+(\w+)'),
        ],
        'rust': [
            re.compile(r'(?:pub(?:\([^)]*\))?\s+)?struct\s+(\w+)'),
            re.compile(r'(?:pub(?:\([^)]*\))?\s+)?trait\s+(\w+)'),
            re.compile(r'(?:pub(?:\([^)]*\))?\s+)?enum\s+(\w+)'),
            re.compile(r'impl(?:<[^>]*>)?\s+(\w+)'),
        ],
    }

    _FUNCTION_PATTERNS = {
        'javascript': [
            re.compile(r'(?:export\s+)?(?:async\s+)?function\s+(\w+)\s*\(([^)]*)\)'),
            re.compile(r'(?:export\s+)?(?:const|let|var)\s+(\w+)\s*=\s*(?:async\s+)?\(([^)]*)\)\s*=>'),
        ],
        'typescript': [
            re.compile(
                r'(?:export\s+)?(?:async\s+)?function\s+(\w+)\s*(?:<[^>]*>)?\s*\(([^)]*)\)'
                r'(?:\s*:\s*[\w<>\[\]|&\s]+)?'
            ),
            re.compile(
                r'(?:export\s+)?(?:const|let|var)\s+(\w+)\s*=\s*(?:async\s+)?'
                r'(?:<[^>]*>)?\s*\(([^)]*)\)\s*(?::\s*[\w<>\[\]|&\s]+)?\s*=>'
            ),
        ],
        'java': [
            re.compile(
                r'(?:public|protected|private)\s+(?:static\s+)?(?:final\s+)?'
                r'(?:synchronized\s+)?(?:[\w<>\[\]]+)\s+(\w+)\s*\(([^)]*)\)'
            ),
        ],
        'go': [
            re.compile(r'func\s+(\w+)\s*\(([^)]*)\)'),
            re.compile(r'func\s+\([^)]+\)\s+(\w+)\s*\(([^)]*)\)'),
        ],
        'c': [
            re.compile(r'^[\w][\w\s*]+\b(\w+)\s*\(([^)]*)\)\s*\{', re.MULTILINE),
        ],
        'cpp': [
            re.compile(
                r'^[\w][\w\s*:&<>,~]+\b(\w+)\s*\(([^)]*)\)\s*'
                r'(?:const\s*)?(?:override\s*)?(?:noexcept\s*)?\{',
                re.MULTILINE
            ),
        ],
        'rust': [
            re.compile(
                r'(?:pub(?:\([^)]*\))?\s+)?(?:async\s+)?fn\s+(\w+)\s*'
                r'(?:<[^>]*>)?\s*\(([^)]*)\)'
            ),
        ],
    }

    @staticmethod
    def analyze(content: str, file_path: str, language: str) -> Dict[str, Any]:
        """Analyze source code using regex patterns for the specified language."""
        classes = _RegexCodeAnalyzer._extract_classes(content, language)
        functions = _RegexCodeAnalyzer._extract_functions(content, language)
        return {"classes": classes, "functions": functions}

    @staticmethod
    def _extract_classes(content: str, language: str) -> List[Dict[str, Any]]:
        """Extract class, struct, interface, or trait definitions."""
        patterns = _RegexCodeAnalyzer._CLASS_PATTERNS.get(language, [])
        classes = []
        seen_names = set()

        for pattern in patterns:
            for match in pattern.finditer(content):
                name = match.group(1)
                if name in seen_names or name in _RegexCodeAnalyzer._KEYWORDS:
                    continue
                seen_names.add(name)

                base_classes = []
                if match.lastindex >= 2 and match.group(2):
                    bases = match.group(2).strip()
                    base_classes = [b.strip() for b in bases.split(',') if b.strip()]

                docstring = _RegexCodeAnalyzer._extract_nearby_comment(
                    content, match.start()
                )

                classes.append({
                    "name": name,
                    "base_classes": base_classes,
                    "docstring": docstring[:200] if docstring else None,
                    "methods": []
                })

        return classes

    @staticmethod
    def _extract_functions(content: str, language: str) -> List[Dict[str, Any]]:
        """Extract function and method definitions."""
        patterns = _RegexCodeAnalyzer._FUNCTION_PATTERNS.get(language, [])
        functions = []
        seen_names = set()

        for pattern in patterns:
            for match in pattern.finditer(content):
                name = match.group(1)
                if name in seen_names or name in _RegexCodeAnalyzer._KEYWORDS:
                    continue
                seen_names.add(name)

                params = []
                if match.lastindex >= 2 and match.group(2):
                    raw_params = match.group(2).strip()
                    if raw_params:
                        params = [p.strip() for p in raw_params.split(',') if p.strip()]
                        params = params[:10]

                return_type = None
                if match.lastindex >= 3 and match.group(3):
                    return_type = match.group(3).strip()

                docstring = _RegexCodeAnalyzer._extract_nearby_comment(
                    content, match.start()
                )

                # Detect async qualifier from line context
                is_async = False
                line_start = content.rfind('\n', 0, match.start()) + 1
                line_prefix = content[line_start:match.start()]
                if 'async' in line_prefix:
                    is_async = True

                functions.append({
                    "name": name,
                    "parameters": params,
                    "return_type": return_type,
                    "docstring": docstring[:200] if docstring else None,
                    "is_async": is_async,
                    "decorators": []
                })

        return functions

    @staticmethod
    def _extract_nearby_comment(content: str, pos: int) -> Optional[str]:
        """Extract documentation comment immediately preceding the given position."""
        preceding = content[max(0, pos - 500):pos].rstrip()

        # JSDoc / JavaDoc / Doxygen block comment: /** ... */
        block_match = re.search(r'/\*\*\s*(.*?)\*/', preceding, re.DOTALL)
        if block_match:
            comment = block_match.group(1).strip()
            lines = [line.strip().lstrip('* ').strip() for line in comment.split('\n')]
            return ' '.join(line for line in lines if line)[:200]

        # Line comments (// or ///)
        lines = preceding.split('\n')
        comment_lines = []
        for line in reversed(lines):
            stripped = line.strip()
            if stripped.startswith('//'):
                comment_lines.insert(0, stripped.lstrip('/').strip())
            elif stripped == '':
                continue
            else:
                break

        if comment_lines:
            return ' '.join(comment_lines)[:200]

        return None


class _OfficeReader:
    """Extract text content from Office documents (PDF, PPT, Word)."""

    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.pptx', '.ppt'}

    @staticmethod
    def is_supported(file_path: str) -> bool:
        """Check if file type is supported."""
        ext = os.path.splitext(file_path)[1].lower()
        return ext in _OfficeReader.SUPPORTED_EXTENSIONS

    @staticmethod
    def get_file_type(file_path: str) -> str:
        """Get human-readable file type."""
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {
            '.pdf': 'PDF Document',
            '.docx': 'Word Document',
            '.doc': 'Word Document (Legacy)',
            '.pptx': 'PowerPoint Presentation',
            '.ppt': 'PowerPoint Presentation (Legacy)'
        }
        return type_map.get(ext, 'Unknown')

    @staticmethod
    def extract_text(file_path: str, max_chars: int = 50000) -> str:
        """
        Extract text content from supported office documents.
        
        Args:
            file_path: Path to the office document
            max_chars: Maximum characters to extract
            
        Returns:
            Extracted text content
            
        Raises:
            ValueError: If file type not supported
            ImportError: If required library not installed
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.pdf':
            return _OfficeReader._extract_pdf(file_path, max_chars)
        elif ext in ('.docx', '.doc'):
            return _OfficeReader._extract_word(file_path, max_chars)
        elif ext in ('.pptx', '.ppt'):
            return _OfficeReader._extract_ppt(file_path, max_chars)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    @staticmethod
    def _extract_pdf(file_path: str, max_chars: int) -> str:
        """Extract text from PDF file."""
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            raise ImportError(
                "PyPDF2 is required for PDF extraction. "
                "Install with: pip install PyPDF2"
            )
        
        try:
            reader = PdfReader(file_path)
            text_parts = []
            total_chars = 0
            
            for page in reader.pages:
                page_text = page.extract_text() or ""
                if total_chars + len(page_text) > max_chars:
                    remaining = max_chars - total_chars
                    text_parts.append(page_text[:remaining])
                    break
                text_parts.append(page_text)
                total_chars += len(page_text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract PDF text: {e}")
            raise

    @staticmethod
    def _extract_word(file_path: str, max_chars: int) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for Word extraction. "
                "Install with: pip install python-docx"
            )
        
        try:
            doc = Document(file_path)
            text_parts = []
            total_chars = 0
            
            for para in doc.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    continue
                if total_chars + len(para_text) > max_chars:
                    remaining = max_chars - total_chars
                    text_parts.append(para_text[:remaining])
                    break
                text_parts.append(para_text)
                total_chars += len(para_text)
            
            # Also extract text from tables
            for table in doc.tables:
                if total_chars >= max_chars:
                    break
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if total_chars + len(row_text) > max_chars:
                        break
                    text_parts.append(row_text)
                    total_chars += len(row_text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract Word text: {e}")
            raise

    @staticmethod
    def _extract_ppt(file_path: str, max_chars: int) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint extraction. "
                "Install with: pip install python-pptx"
            )
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            total_chars = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_texts.append(para_text)
                
                slide_content = "\n".join(slide_texts)
                if total_chars + len(slide_content) > max_chars:
                    remaining = max_chars - total_chars
                    text_parts.append(slide_content[:remaining])
                    break
                text_parts.append(slide_content)
                total_chars += len(slide_content)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to extract PowerPoint text: {e}")
            raise